#!/usr/bin/env python3
"""
Slide Generation Script for GitHub Actions
"""
import os
import json
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np
from PIL import Image
import re
import textwrap

class SlideGenerator:
    def __init__(self, task_id, template_type='professional', theme='modern', language='ja'):
        self.task_id = task_id
        self.template_type = template_type
        self.theme = theme
        self.language = language
        self.output_dir = f"outputs/slides/{task_id}"
        
        # Theme colors
        self.themes = {
            'modern': {
                'primary': RGBColor(41, 128, 185),      # Blue
                'secondary': RGBColor(52, 73, 94),      # Dark Blue
                'accent': RGBColor(231, 76, 60),        # Red
                'background': RGBColor(255, 255, 255),  # White
                'text': RGBColor(44, 62, 80)            # Dark Gray
            }
        }
        
        self.current_theme = self.themes.get(theme, self.themes['modern'])
        self.presentation = Presentation()
        self.slide_layouts = self.presentation.slide_layouts
        
        # Content data
        self.content_data = []
        self.pdf_documents = []
        self.pdf_summary = {}
        self.processed_text = ""
        self.slide_data = []
        
    def load_data_sources(self, data_source):
        """Load data from various sources"""
        sources_loaded = []
        
        # Load scraping data
        scraping_file = f"outputs/scraping/{self.task_id}/detailed_data.json"
        if os.path.exists(scraping_file):
            try:
                with open(scraping_file, 'r', encoding='utf-8') as f:
                    scraping_data = json.load(f)
                
                for page in scraping_data:
                    content = page.get('content', {})
                    self.processed_text += content.get('full_text', '') + "\n"
                    
                    # Extract structured content
                    self.content_data.append({
                        'type': 'web_page',
                        'title': page.get('title', ''),
                        'url': page.get('url', ''),
                        'headings': content.get('headings', []),
                        'paragraphs': content.get('paragraphs', []),
                        'lists': content.get('lists', []),
                        'tables': content.get('tables', [])
                    })
                
                sources_loaded.append('scraping')
                print(f"Loaded scraping data: {len(scraping_data)} pages")
            except Exception as e:
                print(f"Error loading scraping data: {e}")

        # Load PDF processing data
        pdf_results_file = f"outputs/pdf-processing/{self.task_id}/processing_results.json"
        if os.path.exists(pdf_results_file):
            try:
                with open(pdf_results_file, 'r', encoding='utf-8') as f:
                    pdf_data = json.load(f)

                if isinstance(pdf_data, dict) and isinstance(pdf_data.get('summary'), dict):
                    self.pdf_summary = pdf_data['summary']

                if isinstance(pdf_data, list):
                    pdf_entries = pdf_data
                elif isinstance(pdf_data, dict):
                    if isinstance(pdf_data.get('results'), list):
                        pdf_entries = pdf_data['results']
                    elif isinstance(pdf_data.get('documents'), list):
                        pdf_entries = pdf_data['documents']
                    else:
                        pdf_entries = [pdf_data]
                else:
                    pdf_entries = []

                loaded_documents = 0
                for entry in pdf_entries:
                    if not isinstance(entry, dict):
                        continue

                    status = entry.get('status', 'unknown')
                    if status not in ('completed', 'success', 'succeeded'):
                        # Skip documents that did not finish successfully but keep metadata for reference
                        print(f"Skipping PDF entry due to status '{status}': {entry.get('file', 'unknown')}")
                        continue

                    metadata = entry.get('metadata') or {}
                    filename = metadata.get('filename') or os.path.basename(entry.get('file', '')) or 'PDF Document'

                    text_segments = []
                    for key in ('text', 'ocr_text'):
                        value = entry.get(key)
                        if isinstance(value, str) and value.strip():
                            text_segments.append(value.strip())

                    combined_text = "\n".join(text_segments)
                    if combined_text:
                        self.processed_text += combined_text + "\n"

                    images = entry.get('images') if isinstance(entry.get('images'), list) else []
                    tables = entry.get('tables') if isinstance(entry.get('tables'), list) else []

                    document_info = {
                        'type': 'pdf_document',
                        'title': filename,
                        'file': entry.get('file', ''),
                        'metadata': metadata,
                        'page_count': entry.get('page_count'),
                        'tables': tables,
                        'images': images,
                        'text': combined_text,
                        'text_excerpt': combined_text[:500] if combined_text else '',
                        'status': status
                    }

                    self.content_data.append(document_info)
                    self.pdf_documents.append(document_info)
                    loaded_documents += 1

                if loaded_documents:
                    sources_loaded.append('pdf_processing')
                    print(f"Loaded PDF processing data: {loaded_documents} documents")
            except Exception as e:
                print(f"Error loading PDF processing data: {e}")
        else:
            if data_source not in ('manual', None):
                print(f"No PDF processing results found at {pdf_results_file}")

        return sources_loaded

    def analyze_content(self, max_slides=10):
        """Analyze content and create slide structure"""
        print(f"Analyzing content for {max_slides} slides...")

        # Basic content analysis
        words = self.processed_text.split()

        # Extract key topics (simple keyword extraction)
        word_freq = {}
        for word in words:
            word_clean = re.sub(r'[^\w]', '', word.lower())
            if len(word_clean) > 3:  # Skip short words
                word_freq[word_clean] = word_freq.get(word_clean, 0) + 1

        # Get top keywords
        top_keywords = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:20]

        # Generate slide structure
        slides = []

        # Title slide
        slides.append({
            'type': 'title',
            'title': 'データ分析結果' if self.language == 'ja' else 'Data Analysis Results',
            'subtitle': f'Task ID: {self.task_id}',
            'content': []
        })

        # Overview slide
        slides.append({
            'type': 'overview',
            'title': '概要' if self.language == 'ja' else 'Overview',
            'content': [
                f"処理されたデータソース: {len(self.content_data)}" if self.language == 'ja' else f"Processed data sources: {len(self.content_data)}",
                f"総文字数: {len(self.processed_text):,}" if self.language == 'ja' else f"Total characters: {len(self.processed_text):,}",
                f"主要キーワード数: {len(top_keywords)}" if self.language == 'ja' else f"Key keywords: {len(top_keywords)}"
            ]
        })

        pdf_docs = self.pdf_documents if hasattr(self, 'pdf_documents') else []

        if pdf_docs:
            total_pdf_chars = sum(len(doc.get('text', '') or '') for doc in pdf_docs)
            total_pages = sum(doc.get('page_count') or 0 for doc in pdf_docs)

            pdf_overview_line = (
                f"PDFドキュメント数: {len(pdf_docs)} / 総ページ数: {total_pages}" if self.language == 'ja'
                else f"PDF documents: {len(pdf_docs)} / Total pages: {total_pages}"
            )
            slides[1]['content'].append(pdf_overview_line)

            def count_images(image_entries):
                if not isinstance(image_entries, list):
                    return 0
                total = 0
                for item in image_entries:
                    if isinstance(item, dict):
                        if isinstance(item.get('count'), int):
                            total += item['count']
                        elif 'image_index' in item:
                            total += 1
                return total

            total_tables = sum(len(doc.get('tables') or []) for doc in pdf_docs)
            total_images = sum(count_images(doc.get('images')) for doc in pdf_docs)

            slides.append({
                'type': 'content',
                'title': 'PDF分析ハイライト' if self.language == 'ja' else 'PDF Analysis Highlights',
                'content': [
                    f"処理したPDF: {len(pdf_docs)}件" if self.language == 'ja' else f"Processed PDFs: {len(pdf_docs)}",
                    f"総ページ数: {total_pages}" if self.language == 'ja' else f"Total pages: {total_pages}",
                    f"抽出されたテキスト文字数: {total_pdf_chars:,}" if self.language == 'ja' else f"Extracted text characters: {total_pdf_chars:,}",
                    f"抽出されたテーブル数: {total_tables}" if self.language == 'ja' else f"Extracted tables: {total_tables}",
                    f"検出された画像数: {total_images}" if self.language == 'ja' else f"Detected images: {total_images}"
                ]
            })

            if self.pdf_summary:
                summary_slide_content = []
                for key, value in self.pdf_summary.items():
                    summary_slide_content.append(f"{key}: {value}")

                if summary_slide_content:
                    slides.append({
                        'type': 'content',
                        'title': 'PDFサマリー' if self.language == 'ja' else 'PDF Summary',
                        'content': summary_slide_content
                    })

            remaining_slots = max_slides - len(slides)
            if remaining_slots > 0:
                for document in pdf_docs[:remaining_slots]:
                    doc_content = []
                    metadata = document.get('metadata') or {}

                    if metadata.get('filename'):
                        doc_content.append(
                            f"ファイル名: {metadata['filename']}" if self.language == 'ja' else f"Filename: {metadata['filename']}"
                        )

                    if metadata.get('size_bytes'):
                        size_bytes = metadata['size_bytes']
                        doc_content.append(
                            f"ファイルサイズ: {size_bytes:,} バイト" if self.language == 'ja' else f"File size: {size_bytes:,} bytes"
                        )

                    page_count = document.get('page_count')
                    if page_count:
                        doc_content.append(
                            f"ページ数: {page_count}" if self.language == 'ja' else f"Pages: {page_count}"
                        )

                    tables = document.get('tables') or []
                    if tables:
                        doc_content.append(
                            f"テーブル検出数: {len(tables)}" if self.language == 'ja' else f"Tables detected: {len(tables)}"
                        )

                    excerpt = document.get('text_excerpt')
                    if excerpt:
                        doc_content.append(textwrap.shorten(excerpt, width=180, placeholder='...'))

                    if not doc_content:
                        continue

                    slides.append({
                        'type': 'content',
                        'title': document.get('title') or (
                            'PDFドキュメント' if self.language == 'ja' else 'PDF Document'
                        ),
                        'content': doc_content
                    })

        self.slide_data = slides[:max_slides]
        return self.slide_data
    
    def add_title_slide(self, slide_data):
        """Add title slide"""
        slide_layout = self.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = slide_data['title']
        subtitle.text = slide_data['subtitle']
        
        # Style title
        title_font = title.text_frame.paragraphs[0].font
        title_font.size = Pt(44)
        title_font.color.rgb = self.current_theme['primary']
        title_font.bold = True
        
        # Style subtitle
        subtitle_font = subtitle.text_frame.paragraphs[0].font
        subtitle_font.size = Pt(24)
        subtitle_font.color.rgb = self.current_theme['secondary']
    
    def add_content_slide(self, slide_data):
        """Add content slide"""
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = slide_data['title']
        
        # Add content
        text_frame = content.text_frame
        text_frame.clear()
        
        for item in slide_data['content']:
            if item.strip():
                p = text_frame.paragraphs[0] if len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text else text_frame.add_paragraph()
                p.text = str(item)[:200] + "..." if len(str(item)) > 200 else str(item)  # Limit text length
                p.font.size = Pt(18)
                p.font.color.rgb = self.current_theme['text']
        
        # Style title
        title_font = title.text_frame.paragraphs[0].font
        title_font.size = Pt(32)
        title_font.color.rgb = self.current_theme['primary']
        title_font.bold = True
    
    def generate_presentation(self, max_slides=10):
        """Generate complete presentation"""
        print(f"Generating presentation with {max_slides} slides...")
        
        # Analyze content and create slide structure
        slides = self.analyze_content(max_slides)
        
        # Generate slides
        for slide_data in slides:
            slide_type = slide_data['type']
            
            if slide_type == 'title':
                self.add_title_slide(slide_data)
            elif slide_type in ['overview', 'content', 'conclusion']:
                self.add_content_slide(slide_data)
        
        # Save presentation
        output_filename = f"presentation_{self.task_id}.pptx"
        output_path = os.path.join(self.output_dir, output_filename)
        self.presentation.save(output_path)
        
        # Generate summary
        summary = {
            'task_id': self.task_id,
            'presentation_file': output_filename,
            'slide_count': len(slides),
            'template_type': self.template_type,
            'theme': self.theme,
            'language': self.language,
            'data_sources': len(self.content_data),
            'total_text_length': len(self.processed_text),
            'generation_time': datetime.now().isoformat()
        }
        
        # Save summary
        summary_file = os.path.join(self.output_dir, 'summary.json')
        with open(summary_file, 'w', encoding='utf-8') as f:
            json.dump(summary, f, ensure_ascii=False, indent=2)
        
        return output_path, summary

# Execute slide generation
def main():
    import argparse
    parser = argparse.ArgumentParser(description='Generate presentation slides')
    parser.add_argument('--task-id', required=True)
    parser.add_argument('--template-type', default='professional')
    parser.add_argument('--theme', default='modern')
    parser.add_argument('--language', default='ja')
    parser.add_argument('--data-source', default='auto')
    parser.add_argument('--slide-count', type=int, default=10)
    
    args = parser.parse_args()
    
    generator = SlideGenerator(
        task_id=args.task_id,
        template_type=args.template_type,
        theme=args.theme,
        language=args.language
    )
    
    # Load data sources
    sources = generator.load_data_sources(args.data_source)
    
    if not sources and args.data_source != 'manual':
        print("No data sources found, creating sample presentation...")
        generator.processed_text = "Sample content for presentation generation demonstration."
        generator.content_data = [
            {
                'type': 'sample',
                'title': 'Sample Data',
                'content': 'This is a sample presentation generated without input data.'
            }
        ]
    
    # Generate presentation
    try:
        output_file, summary = generator.generate_presentation(args.slide_count)
        
        print(f"Slide generation completed successfully!")
        print(f"Output file: {output_file}")
        print(f"Summary: {json.dumps(summary, indent=2)}")
        
        # Set GitHub Action outputs
        gh_out = os.getenv('GITHUB_OUTPUT')
        if gh_out:
            with open(gh_out, 'a') as f:
                f.write('status=success\n')
                f.write(f'output_file={output_file}\n')
                f.write(f'slide_summary={json.dumps(summary)}\n')
        
    except Exception as e:
        print(f"Slide generation failed: {e}")
        import traceback
        traceback.print_exc()
        gh_out = os.getenv('GITHUB_OUTPUT')
        if gh_out:
            with open(gh_out, 'a') as f:
                f.write('status=failed\n')

if __name__ == "__main__":
    main()